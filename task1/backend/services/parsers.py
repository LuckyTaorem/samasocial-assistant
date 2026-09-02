import io
import requests
from pptx import Presentation
from youtube_transcript_api import YouTubeTranscriptApi
from bs4 import BeautifulSoup
from urllib.parse import urlparse, parse_qs
from fastapi import HTTPException
import PyPDF2

def parse_pdf(file_bytes):
    reader = PyPDF2.PdfReader(io.BytesIO(file_bytes))
    docs = []
    
    # Loop through each page independently to preserve correct page numbering
    for i, page in enumerate(reader.pages):
        page_text = page.extract_text()
        
        if page_text and page_text.strip():
            docs.append({
                "text": page_text.strip(),
                "metadata": {
                    "page": i + 1,
                    "type": "pdf"
                }
            })
            
    return docs

def parse_pptx(file_bytes):
    prs = Presentation(io.BytesIO(file_bytes))
    print(f"Total slides found in presentation: {len(prs.slides)}")
    docs = []
    
    # Loop through each slide independently to preserve correct slide numbering
    for i, slide in enumerate(prs.slides):
        slide_text = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        slide_text += run.text + " "
                    slide_text += "\n"
        
        slide_text = slide_text.strip()
        if slide_text:
            docs.append({
                "text": slide_text,
                "metadata": {
                    "slide": i + 1,
                    "type": "presentation"
                }
            })
    return docs

def parse_youtube(url: str):
    query = urlparse(url)
    video_id = query.path[1:] if query.hostname == 'youtu.be' else parse_qs(query.query).get('v', [None])[0]
    
    if not video_id:
        raise HTTPException(status_code=400, detail="Invalid YouTube URL")
        
    try:
        # 1. Instantiate the API (Required in the new library version)
        ytt_api = YouTubeTranscriptApi()
        
        # 2. Use the new .list() method to find all available captions
        transcript_list = ytt_api.list(video_id)
        
        try:
            # Try to grab English (manual or auto-generated)
            transcript = transcript_list.find_transcript(['en', 'en-US', 'en-GB', 'en-IN']).fetch()
        except:
            # Fallback: Grab whatever language exists and dynamically translate it to English!
            first_transcript = next(iter(transcript_list))
            transcript = first_transcript.translate('en').fetch()
            
    except Exception as e:
        err_type = str(type(e))
        if "RequestBlocked" in err_type or "IPBlocked" in err_type:
            raise HTTPException(
                status_code=400, 
                detail="YouTube blocked this server's IP address. Cloud hosting providers are often blocked by YouTube's anti-bot systems."
            )
        elif "TranscriptsDisabled" in err_type or "NoTranscript" in err_type:
            raise HTTPException(
                status_code=400, 
                detail="This YouTube video does not have any manual or auto-generated captions enabled."
            )
        else:
            raise HTTPException(
                status_code=400, 
                detail=f"Failed to fetch YouTube transcript. Error: {str(e)}"
            )
    
    docs = []
    current_text_chunks = []
    current_length = 0
    current_timestamp = 0
    
    # 3. Use the new object syntax (entry.text / entry.start) with the fast array-joiner
    for entry in transcript:
        if not current_text_chunks:
            current_timestamp = round(entry.start)
            
        text_part = entry.text
        current_text_chunks.append(text_part)
        current_length += len(text_part) + 1 
        
        if current_length > 800:
            docs.append({
                "text": " ".join(current_text_chunks).strip(),
                "metadata": {"timestamp": current_timestamp}
            })
            current_text_chunks = []
            current_length = 0
            
    if current_text_chunks:
        docs.append({
            "text": " ".join(current_text_chunks).strip(),
            "metadata": {"timestamp": current_timestamp}
        })
        
    return docs

def parse_webpage(url: str):
    response = requests.get(url, headers={"User-Agent": "Mozilla/5.0"})
    soup = BeautifulSoup(response.text, 'html.parser')
    
    for script in soup(["script", "style", "nav", "footer"]):
        script.extract()
        
    text = soup.get_text(separator='\n', strip=True)
    
    return [{"text": text, "metadata": {"source_name": url}}]